"""JSON slide data -> python-pptx PPTX bytes"""
import base64
import io
import math
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from gradient import parse_gradient, css_color_to_rgb, css_color_to_rgba, css_color_to_hex
from text_runs import html_to_text_runs

PX_TO_INCH = 1 / 96
PX_TO_EMU = 914400 / 96  # 9525


def build_pptx(pages: dict, default_canvas_size: dict, fonts: list = None) -> bytes:
    prs = Presentation()

    first_page = next(iter(pages.values()), None)
    cs = (first_page or {}).get('canvasSize', default_canvas_size) or default_canvas_size
    slide_w = cs['w'] * PX_TO_EMU
    slide_h = cs['h'] * PX_TO_EMU
    prs.slide_width = int(slide_w)
    prs.slide_height = int(slide_h)

    # ── Font resolution (before slide building, so name map is available) ──
    font_name_map = {}
    font_records = []
    if fonts:
        try:
            from font_embedder import resolve_fonts, build_font_name_map
            font_records = resolve_fonts(fonts, pages=pages)
            font_name_map = build_font_name_map(font_records)
            print(f'Font embed: resolved {len(font_records)} fonts, {len(font_name_map)} mappings')
        except Exception as e:
            print(f'Font embed: resolution failed: {e}')

    sorted_keys = sorted(pages.keys(), key=_page_sort_key)

    # ── Phase 1: Analyze backgrounds across all pages ──
    page_bg_info = {}  # key -> { 'bg_elements': [...], 'content_elements': [...], 'fingerprint': str }
    fingerprint_counts = {}  # fingerprint -> count

    for key in sorted_keys:
        page = pages[key]
        page_cs = page.get('canvasSize') or cs
        elements = sorted(page.get('elements', []), key=lambda e: e.get('zIndex', 0))
        bg_els, content_els = _separate_background_elements(elements, page_cs)
        fp = _background_fingerprint(bg_els)
        page_bg_info[key] = {
            'bg_elements': bg_els,
            'content_elements': content_els,
            'fingerprint': fp,
            'page_cs': page_cs,
        }
        fingerprint_counts[fp] = fingerprint_counts.get(fp, 0) + 1

    # ── Phase 2: Find the most common background → set as master ──
    most_common_fp = max(fingerprint_counts, key=fingerprint_counts.get) if fingerprint_counts else ''
    # Find a representative page with the most common background
    master_bg_els = []
    for key in sorted_keys:
        if page_bg_info[key]['fingerprint'] == most_common_fp:
            master_bg_els = page_bg_info[key]['bg_elements']
            break

    # Set master slide background
    master = prs.slide_masters[0]
    _apply_background_to_element(master, master_bg_els)
    # Also clear the blank layout's background so it inherits from master
    blank_layout = prs.slide_layouts[6]
    _clear_layout_background(blank_layout)

    # ── Phase 3: Create slides ──
    for key in sorted_keys:
        info = page_bg_info[key]
        page_cs = info['page_cs']
        slide = prs.slides.add_slide(blank_layout)

        # If this page's background differs from master, apply slide-level override
        if info['fingerprint'] != most_common_fp:
            _apply_background_to_element(slide, info['bg_elements'])
            # Add decorative bg layers (grid patterns, etc.) as locked shapes
            for el in info['bg_elements']:
                if _is_decorative_bg_layer(el):
                    try:
                        _add_element(slide, el, page_cs, font_name_map)
                    except Exception:
                        pass
        else:
            # Master handles the main background, but add decorative overlays
            for el in info['bg_elements']:
                if _is_decorative_bg_layer(el):
                    try:
                        _add_element(slide, el, page_cs, font_name_map)
                    except Exception:
                        pass

        # Add content elements
        for el in info['content_elements']:
            try:
                _add_element(slide, el, page_cs, font_name_map)
            except Exception as e:
                print(f'PPT export: element {el.get("id")} skipped: {e}')

    # ── Embed fonts into PPTX package ──
    if font_records:
        try:
            from font_embedder import embed_fonts_in_pptx
            embed_fonts_in_pptx(prs, font_records)
            print(f'Font embed: embedded {len(font_records)} fonts')
        except Exception as e:
            print(f'Font embed: embedding failed: {e}')

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _page_sort_key(key: str):
    parts = key.split('-')
    return tuple(int(p) for p in parts)


def _px(v):
    return int(v * PX_TO_EMU)


def _separate_background_elements(elements: list, cs: dict) -> tuple:
    """Separate full-slide background elements from content elements.
    Background = full-size (>=95% of canvas), locked or z<=2, shape/image type."""
    bg_elements = []
    content_elements = []
    cw, ch = cs.get('w', 1280), cs.get('h', 800)

    for el in elements:
        w = el.get('width', 0)
        h = el.get('height', 0)
        is_fullsize = w >= cw * 0.95 and h >= ch * 0.95
        is_bg_type = el.get('type') in ('shape', 'image')
        z = el.get('zIndex', 0)
        is_locked = el.get('locked', False)
        if is_fullsize and is_bg_type and (is_locked or z <= 2):
            bg_elements.append(el)
        else:
            content_elements.append(el)

    return bg_elements, content_elements


def _background_fingerprint(bg_elements: list) -> str:
    """Create a fingerprint string for a set of background elements.
    Used to group pages with identical backgrounds."""
    parts = []
    for el in bg_elements:
        s = el.get('styles', {})
        bg_color = s.get('backgroundColor', '')
        bg_image = s.get('backgroundImage', 'none')
        # Only include the primary visual properties
        if bg_color and bg_color not in ('rgba(0, 0, 0, 0)', 'transparent'):
            parts.append(f'c:{bg_color}')
        if bg_image and bg_image != 'none' and not _is_subtle_gradient(bg_image):
            parts.append(f'g:{bg_image[:100]}')
    return '|'.join(parts) or 'empty'


def _is_decorative_bg_layer(el: dict) -> bool:
    """Check if a background element is a decorative overlay (grid, subtle pattern)
    that can't be represented as a simple slide background fill."""
    s = el.get('styles', {})
    bg_image = s.get('backgroundImage', 'none')
    bg_color = s.get('backgroundColor', '')

    # Has a gradient/image that IS subtle → decorative overlay
    if bg_image and bg_image != 'none' and _is_subtle_gradient(bg_image):
        return True
    # repeating-* gradients are always decorative patterns
    if bg_image and bg_image.startswith('repeating-'):
        return True
    return False


def _extract_bg_fill(bg_elements: list):
    """Extract the primary background fill from elements.
    Returns ('solid', rgba) or ('gradient', grad_dict) or None."""
    solid_bg = None
    gradient_bg = None

    for el in bg_elements:
        s = el.get('styles', {})
        bg_color = s.get('backgroundColor', '')
        bg_image = s.get('backgroundImage', 'none')

        if bg_image and bg_image != 'none' and not _is_subtle_gradient(bg_image):
            if not bg_image.startswith('repeating-'):
                grad = parse_gradient(bg_image)
                if grad['type'] != 'none' and len(grad['stops']) >= 2:
                    gradient_bg = grad

        if bg_color and bg_color not in ('rgba(0, 0, 0, 0)', 'transparent', ''):
            rgba = css_color_to_rgba(bg_color)
            if rgba and rgba[3] > 0.5:
                solid_bg = rgba

    if gradient_bg:
        return ('gradient', gradient_bg)
    elif solid_bg:
        return ('solid', solid_bg)
    return None


def _apply_bg_fill_xml(parent_el, fill_info, tag_name='p:bgPr'):
    """Apply background fill to a OOXML element (master, layout, or slide).
    parent_el should be the <p:bg> element."""
    # Remove existing bgPr/bgRef
    for tag in ('p:bgPr', 'p:bgRef'):
        existing = parent_el.find(qn(tag))
        if existing is not None:
            parent_el.remove(existing)

    bg_pr = parent_el.makeelement(qn('p:bgPr'), {})

    if fill_info is None:
        # White default
        solid_fill = bg_pr.makeelement(qn('a:solidFill'), {})
        srgb = solid_fill.makeelement(qn('a:srgbClr'), {'val': 'FFFFFF'})
        solid_fill.append(srgb)
        bg_pr.append(solid_fill)
    elif fill_info[0] == 'solid':
        r, g, b, a = fill_info[1]
        solid_fill = bg_pr.makeelement(qn('a:solidFill'), {})
        srgb = solid_fill.makeelement(qn('a:srgbClr'), {'val': f'{r:02X}{g:02X}{b:02X}'})
        if a < 0.95:
            alpha_el = srgb.makeelement(qn('a:alpha'), {'val': str(int(a * 100000))})
            srgb.append(alpha_el)
        solid_fill.append(srgb)
        bg_pr.append(solid_fill)
    elif fill_info[0] == 'gradient':
        grad = fill_info[1]
        grad_fill = bg_pr.makeelement(qn('a:gradFill'), {})
        gs_lst = grad_fill.makeelement(qn('a:gsLst'), {})
        for stop in grad['stops']:
            gs = gs_lst.makeelement(qn('a:gs'), {'pos': str(int(stop['position'] * 1000))})
            rgb = css_color_to_rgb(stop['color'])
            if rgb:
                srgb = gs.makeelement(qn('a:srgbClr'), {'val': f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'})
                gs.append(srgb)
                gs_lst.append(gs)
        grad_fill.append(gs_lst)
        if grad['type'] == 'linear':
            ooxml_angle = int(((grad['angle'] + 270) % 360) * 60000)
            lin = grad_fill.makeelement(qn('a:lin'), {'ang': str(ooxml_angle), 'scaled': '0'})
            grad_fill.append(lin)
        bg_pr.append(grad_fill)

    effect_lst = bg_pr.makeelement(qn('a:effectLst'), {})
    bg_pr.append(effect_lst)
    parent_el.append(bg_pr)


def _apply_background_to_element(target, bg_elements: list):
    """Apply background fill to a slide master, layout, or slide.
    Creates <p:bg> with proper fill from background elements."""
    fill_info = _extract_bg_fill(bg_elements)

    # Force creation of <p:bg> element via python-pptx API
    bg_obj = target.background
    bg_obj.fill.solid()
    bg_obj.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Now find and replace the created <p:bg> content
    # For slide master: element is <p:sldMaster>, child is <p:cSld>
    # For slide: element is <p:sld>, child is <p:cSld>
    cSld = target._element.find(qn('p:cSld'))
    if cSld is None:
        return
    bg_el = cSld.find(qn('p:bg'))
    if bg_el is None:
        return

    _apply_bg_fill_xml(bg_el, fill_info)


def _clear_layout_background(layout):
    """Remove any explicit background from a layout so it inherits from master."""
    cSld = layout._element.find(qn('p:cSld'))
    if cSld is not None:
        bg_el = cSld.find(qn('p:bg'))
        if bg_el is not None:
            cSld.remove(bg_el)


def _clear_theme_style(shape):
    """Remove the default <p:style> element that python-pptx adds to shapes.
    This prevents theme-based effects (shadows from effectRef idx=2)
    and theme fills (fillRef idx=3) from overriding our custom styling."""
    sp = shape._element
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    style_el = sp.find(f'{ns_p}style')
    if style_el is not None:
        sp.remove(style_el)


def _is_background_element(el: dict, cs: dict) -> bool:
    """Check if element is a full-slide background (for locking in PPT)."""
    cw, ch = cs.get('w', 1280), cs.get('h', 800)
    w = el.get('width', 0)
    h = el.get('height', 0)
    is_fullsize = w >= cw * 0.95 and h >= ch * 0.95
    is_locked = el.get('locked', False)
    z = el.get('zIndex', 0)
    return is_fullsize and (is_locked or z <= 2)


def _lock_shape(shape):
    """Lock a shape in PPT so it can't be easily selected (noSelect, noMove, noResize)."""
    try:
        sp = shape._element
        # For shapes (p:sp), lock via cNvSpPr > a:spLocks
        cNv_sp_pr = sp.find(qn('p:nvSpPr'))
        if cNv_sp_pr is not None:
            cNv_tag = cNv_sp_pr.find(qn('p:cNvSpPr'))
            if cNv_tag is None:
                cNv_tag = cNv_sp_pr.makeelement(qn('p:cNvSpPr'), {})
                cNv_sp_pr.append(cNv_tag)
            locks = cNv_tag.find(qn('a:spLocks'))
            if locks is None:
                locks = cNv_tag.makeelement(qn('a:spLocks'), {})
                cNv_tag.append(locks)
            locks.set('noSelect', '1')
            locks.set('noMove', '1')
            locks.set('noResize', '1')
            return

        # For pictures (p:pic), lock via cNvPicPr > a:picLocks
        cNv_pic_pr = sp.find(qn('p:nvPicPr'))
        if cNv_pic_pr is not None:
            cNv_tag = cNv_pic_pr.find(qn('p:cNvPicPr'))
            if cNv_tag is None:
                cNv_tag = cNv_pic_pr.makeelement(qn('p:cNvPicPr'), {})
                cNv_pic_pr.append(cNv_tag)
            locks = cNv_tag.find(qn('a:picLocks'))
            if locks is None:
                locks = cNv_tag.makeelement(qn('a:picLocks'), {})
                cNv_tag.append(locks)
            locks.set('noSelect', '1')
            locks.set('noMove', '1')
            locks.set('noResize', '1')
    except Exception:
        pass


def _add_element(slide, el: dict, cs: dict, font_name_map: dict = None):
    x = _px(el.get('x', 0))
    y = _px(el.get('y', 0))
    w = _px(el.get('width', 0))
    h = _px(el.get('height', 0))
    rotation = el.get('rotation', 0)
    s = el.get('styles', {})
    el_type = el.get('type', '')

    # Track shape count before adding, to lock background elements after
    is_bg = _is_background_element(el, cs)
    shape_count_before = len(slide.shapes)

    if el_type == 'text':
        _add_text(slide, el, x, y, w, h, rotation, font_name_map)
    elif el_type == 'image':
        _add_image(slide, el, x, y, w, h, rotation)
    elif el_type == 'shape':
        _add_shape(slide, el, x, y, w, h, rotation, cs)
    elif el_type == 'svg':
        _add_svg(slide, el, x, y, w, h, rotation)
    elif el_type == 'video':
        _add_video_placeholder(slide, el, x, y, w, h, rotation)

    # Lock background shapes so they can't be accidentally selected in PPT
    if is_bg and len(slide.shapes) > shape_count_before:
        for si in range(shape_count_before, len(slide.shapes)):
            _lock_shape(slide.shapes[si])


def _add_text(slide, el: dict, x, y, w, h, rotation, font_name_map: dict = None):
    s = el.get('styles', {})
    content = el.get('content', '')

    # gradient text detection
    is_gradient_text = (
        s.get('webkitBackgroundClip') == 'text' or s.get('backgroundClip') == 'text'
    ) and s.get('backgroundImage', 'none') != 'none'

    effective_color = s.get('color', '')
    if is_gradient_text:
        grad = parse_gradient(s.get('backgroundImage', ''))
        if grad['stops']:
            effective_color = grad['stops'][0]['color']
    if s.get('webkitTextFillColor') in ('transparent', 'rgba(0, 0, 0, 0)'):
        if not is_gradient_text:
            effective_color = s.get('color', '')

    txbox = slide.shapes.add_textbox(x, y, w, h)
    _clear_theme_style(txbox)
    if rotation:
        txbox.rotation = rotation

    tf = txbox.text_frame
    tf.word_wrap = True

    # Detect circular container (for forced centering)
    br = s.get('borderRadius', '0px')
    is_circle = _is_ellipse(br, w, h) if br and br not in ('0px', '') else False

    # vertical alignment
    has_bg = s.get('backgroundColor', '') not in ('', 'rgba(0, 0, 0, 0)', 'transparent')
    if is_circle:
        # 원형 컨테이너: 항상 수평/수직 중앙 정렬
        _set_anchor(txbox, 'middle')
    elif el.get('merged') or has_bg:
        ai = s.get('alignItems', 'center') if s.get('isFlex') else 'center'
        if ai == 'center':
            _set_anchor(txbox, 'middle')
        elif ai == 'flex-end':
            _set_anchor(txbox, 'bottom')
    else:
        _set_anchor(txbox, 'top')

    # margin (internal padding) - default to 0
    _set_margins(txbox, s.get('padding', '0px'))

    # auto-fit: shrink text to fit if overflow (normAutofit)
    body_props = tf._txBody.find(qn('a:bodyPr'))
    if body_props is not None:
        # Remove default spAutoFit (shape resizes to text)
        for auto_tag in ('a:spAutoFit', 'a:noAutofit', 'a:normAutofit'):
            existing = body_props.find(qn(auto_tag))
            if existing is not None:
                body_props.remove(existing)
        # normAutofit: shrinks font to fit when text overflows, keeps original size when it fits
        norm_auto = body_props.makeelement(qn('a:normAutofit'), {'fontScale': '100000', 'lnSpcReduction': '0'})
        body_props.append(norm_auto)

    # background fill
    bg_result = _parse_solid_fill(s)
    bg_image = s.get('backgroundImage', 'none')
    has_bg_gradient = bg_image != 'none' and not is_gradient_text

    if bg_result:
        bg_rgb, bg_alpha = bg_result
        txbox.fill.solid()
        txbox.fill.fore_color.rgb = bg_rgb
        if bg_alpha < 1.0:
            _set_fill_transparency(txbox, int(bg_alpha * 100000))
    elif has_bg_gradient:
        grad = parse_gradient(bg_image)
        if grad['type'] != 'none' and len(grad['stops']) >= 2:
            _apply_gradient_fill_to_shape(txbox, grad)
        else:
            txbox.fill.background()
    else:
        txbox.fill.background()

    # border (uniform)
    _apply_border(txbox, s)
    # partial borders (1-2 sides) as separate line shapes
    _add_partial_borders(slide, s, x, y, w, h)

    # borderRadius (TextBox도 roundRect 적용)
    br = s.get('borderRadius', '0px')
    if br and br not in ('0px', ''):
        if _is_ellipse(br, w, h):
            sp_pr = txbox._element.spPr
            prst_geom = sp_pr.find(qn('a:prstGeom'))
            if prst_geom is not None:
                prst_geom.set('prst', 'ellipse')
        else:
            _apply_round_corners(txbox, br)

    # shadow
    _apply_shadow(txbox, s)
    _apply_text_shadow(txbox, s)

    # opacity
    _apply_opacity(txbox, s)

    # Resolve font weight from fontWeight + font-variation-settings "wght" axis
    effective_weight = _resolve_font_weight(s)

    # build text runs
    base_styles = {**s, 'color': effective_color, '_effectiveWeight': effective_weight}
    if el.get('isRich') and content:
        runs_data = html_to_text_runs(content, base_styles)
    else:
        opts = {}
        if effective_color:
            opts['color'] = css_color_to_hex(effective_color)
        fs = s.get('fontSize')
        if fs:
            opts['fontSize'] = round(float(str(fs).replace('px', '')) * 0.75)
        ff = s.get('fontFamily')
        if ff:
            opts['fontFace'] = ff.split(',')[0].strip().strip("'\"")
        if effective_weight >= 700:
            opts['bold'] = True
        if s.get('fontStyle') == 'italic':
            opts['italic'] = True
        opts['_weight'] = effective_weight
        runs_data = [{'text': content, 'opts': opts}]

    # text alignment
    if is_circle:
        align = PP_ALIGN.CENTER
    elif s.get('isFlex') and s.get('justifyContent') == 'center':
        align = PP_ALIGN.CENTER
    else:
        align_map = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT, 'left': PP_ALIGN.LEFT}
        align = align_map.get(s.get('textAlign', ''), PP_ALIGN.LEFT)

    # line spacing
    line_spacing = None
    lh = s.get('lineHeight')
    if lh:
        try:
            lh_val = float(lh)
            fs_val = float(str(s.get('fontSize', '16')).replace('px', ''))
            line_spacing = Pt(round(lh_val * fs_val * 0.75))
        except (ValueError, TypeError):
            pass

    _populate_text_frame(tf, runs_data, align, line_spacing, font_name_map)


def _populate_text_frame(tf, runs_data: list, align, line_spacing, font_name_map: dict = None):
    # split runs by newlines into paragraphs
    paragraphs = [[]]
    for run in runs_data:
        text = run['text']
        if '\n' in text:
            parts = text.split('\n')
            for i, part in enumerate(parts):
                if i > 0:
                    paragraphs.append([])
                if part:
                    paragraphs[-1].append({'text': part, 'opts': run['opts']})
        else:
            paragraphs[-1].append(run)

    for pi, para_runs in enumerate(paragraphs):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing

        if not para_runs:
            continue

        for ri, run_data in enumerate(para_runs):
            if pi == 0 and ri == 0:
                r = p.runs[0] if p.runs else p.add_run()
            else:
                r = p.add_run()
            r.text = run_data['text']
            opts = run_data['opts']
            if opts.get('bold'):
                r.font.bold = True
            if opts.get('italic'):
                r.font.italic = True
            if opts.get('underline'):
                r.font.underline = True
            color_hex = opts.get('color')
            if color_hex:
                r.font.color.rgb = RGBColor.from_string(color_hex)
            font_size = opts.get('fontSize')
            if font_size:
                r.font.size = Pt(font_size)
            font_face = opts.get('fontFace')
            if font_face:
                # Map (family, weight) to embedded font typeface if available
                if font_name_map:
                    weight = opts.get('_weight', 700 if opts.get('bold') else 400)
                    font_face = font_name_map.get((font_face, weight), font_face)
                r.font.name = font_face
            # letter-spacing → OOXML spc (in 1/100 pt)
            letter_spacing = opts.get('letterSpacing')
            if letter_spacing is not None:
                rPr = r._r.get_or_add_rPr()
                rPr.set('spc', str(int(letter_spacing * 100)))


def _add_image(slide, el: dict, x, y, w, h, rotation):
    content = el.get('content', '')
    s = el.get('styles', {})

    if content.startswith('data:'):
        m = re.match(r'data:[^;]+;base64,(.+)', content)
        if not m:
            return
        img_bytes = base64.b64decode(m.group(1))
    else:
        return  # external URLs not supported server-side

    img_stream = io.BytesIO(img_bytes)
    pic = slide.shapes.add_picture(img_stream, x, y, w, h)
    if rotation:
        pic.rotation = rotation


def _add_shape(slide, el: dict, x, y, w, h, rotation, cs: dict):
    s = el.get('styles', {})

    bg_image = s.get('backgroundImage', 'none')
    has_gradient = bg_image != 'none'
    # Skip subtle decorative radial gradients (low-alpha fades to transparent)
    if has_gradient and _is_subtle_gradient(bg_image):
        has_gradient = False
    fill_result = _parse_solid_fill(s)
    border_data = _parse_border_data(s)
    partial_borders = _parse_partial_borders(s)
    shadow_data = s.get('boxShadow', '')

    has_any_border = border_data or (partial_borders and not border_data)
    if not has_gradient and not fill_result and not has_any_border and shadow_data in ('', 'none'):
        return

    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.RECTANGLE = 1
    _clear_theme_style(shape)
    if rotation:
        shape.rotation = rotation

    # border radius
    br = s.get('borderRadius', '0px')
    if _is_ellipse(br, w, h):
        sp_pr = shape._element.spPr
        prst_geom = sp_pr.find(qn('a:prstGeom'))
        if prst_geom is not None:
            prst_geom.set('prst', 'ellipse')
    elif br and br != '0px':
        _apply_round_corners(shape, br)

    # fill
    if has_gradient:
        grad = parse_gradient(s.get('backgroundImage', ''))
        if grad['type'] != 'none' and len(grad['stops']) >= 2:
            _apply_gradient_fill_to_shape(shape, grad)
        elif fill_result:
            fill_rgb, fill_alpha = fill_result
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
            if fill_alpha < 1.0:
                _set_fill_transparency(shape, int(fill_alpha * 100000))
        else:
            shape.fill.background()
    elif fill_result:
        fill_rgb, fill_alpha = fill_result
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        if fill_alpha < 1.0:
            _set_fill_transparency(shape, int(fill_alpha * 100000))
    else:
        shape.fill.background()

    # border (uniform)
    _apply_border(shape, s)
    # partial borders (1-2 sides) as separate line shapes
    _add_partial_borders(slide, s, x, y, w, h)

    # shadow
    _apply_shadow(shape, s)

    # opacity
    _apply_opacity(shape, s)


def _add_svg(slide, el: dict, x, y, w, h, rotation):
    content = el.get('content', '')
    if not content:
        return
    # Convert SVG to PNG via cairosvg, then insert as image
    try:
        import cairosvg
        svg_bytes = content.encode('utf-8')
        # Calculate output size in pixels for good quality
        w_px = max(int(w / PX_TO_EMU * 2), 64)  # 2x for retina quality
        h_px = max(int(h / PX_TO_EMU * 2), 64)
        png_bytes = cairosvg.svg2png(bytestring=svg_bytes, output_width=w_px, output_height=h_px)
        img_stream = io.BytesIO(png_bytes)
        pic = slide.shapes.add_picture(img_stream, x, y, w, h)
        if rotation:
            pic.rotation = rotation
    except Exception as e:
        print(f'SVG export failed: {e}')
        # Fallback: try direct SVG insert
        try:
            img_stream = io.BytesIO(content.encode('utf-8'))
            pic = slide.shapes.add_picture(img_stream, x, y, w, h)
            if rotation:
                pic.rotation = rotation
        except Exception:
            pass


def _add_video_placeholder(slide, el: dict, x, y, w, h, rotation):
    txbox = slide.shapes.add_textbox(x, y, w, h)
    _clear_theme_style(txbox)
    if rotation:
        txbox.rotation = rotation
    txbox.fill.solid()
    txbox.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f'\u25b6 video\n{el.get("content", "")}'
    r.font.size = Pt(10)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ── Helpers ──

def _resolve_font_weight(s: dict) -> int:
    """Resolve effective font weight from CSS fontWeight + font-variation-settings 'wght' axis."""
    # 1. Check font-variation-settings for explicit wght value
    fvs = s.get('fontVariationSettings', '')
    if fvs and fvs != 'normal':
        m = re.search(r'["\']wght["\']\s+([\d.]+)', fvs)
        if m:
            return int(float(m.group(1)))

    # 2. Fall back to fontWeight
    fw = s.get('fontWeight', '400')
    if fw == 'bold':
        return 700
    if fw == 'normal':
        return 400
    try:
        return int(fw)
    except (ValueError, TypeError):
        return 400


def _is_subtle_gradient(bg_image: str) -> bool:
    """Skip decorative gradients with very low alpha (glows, glass effects, subtle overlays).
    Catches both radial and linear gradients that fade to transparent."""
    if not bg_image or bg_image == 'none':
        return False
    # Extract all rgba alpha values
    alphas = re.findall(r'rgba\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*,\s*([\d.]+)\s*\)', bg_image)
    if not alphas:
        return False
    max_alpha = max(float(a) for a in alphas)
    # Very low alpha gradients are decorative overlays (glass, glow, etc.)
    if max_alpha < 0.3:
        return True
    # Gradients that fade to fully transparent are often decorative
    min_alpha = min(float(a) for a in alphas)
    if min_alpha < 0.05 and max_alpha < 0.5:
        return True
    return False


def _is_ellipse(border_radius: str, w_emu: int, h_emu: int) -> bool:
    """Detect if borderRadius creates a true circle/ellipse shape.
    Only returns True for roughly square shapes (aspect ratio < 1.5:1)
    where radius >= 45% of the min dimension.
    Wide rounded rectangles (pills/capsules) should use roundRect instead."""
    if not border_radius:
        return False
    br = border_radius.strip()

    # Check aspect ratio first: only square-ish shapes can be ellipses
    # Wide/tall rectangles with borderRadius:50% are pills, not ellipses
    w_px = w_emu / PX_TO_EMU if w_emu else 0
    h_px = h_emu / PX_TO_EMU if h_emu else 0
    if w_px < 1 or h_px < 1:
        return False
    aspect = max(w_px, h_px) / min(w_px, h_px)
    if aspect > 1.5:
        return False  # too elongated → pill shape, use roundRect

    if br in ('50%', '9999px'):
        return True

    # Parse all radius values
    parts = br.split()
    px_values = []
    for p in parts:
        try:
            px_values.append(float(p.replace('px', '').replace('%', '')))
        except (ValueError, TypeError):
            continue
    if not px_values:
        return False
    min_radius_px = min(px_values)
    min_dim_px = min(w_px, h_px)
    # If radius >= ~45% of min dimension, treat as ellipse
    return min_dim_px > 0 and min_radius_px >= min_dim_px * 0.45

def _parse_solid_fill(s: dict):
    """Returns (RGBColor, alpha) or None. alpha < 0.05 → transparent."""
    bg = s.get('backgroundColor', '')
    if bg and bg not in ('rgba(0, 0, 0, 0)', 'transparent'):
        rgba = css_color_to_rgba(bg)
        if rgba:
            r, g, b, a = rgba
            if a < 0.05:
                return None
            return RGBColor(r, g, b), a
    return None


def _parse_border_one(border_str: str) -> dict | None:
    m = re.match(r'([\d.]+)px\s+(\w+)\s+(.+)', border_str)
    if not m:
        return None
    color_str = m.group(3).strip()
    style = m.group(2).lower()  # solid, dashed, dotted, double, etc.
    rgba = css_color_to_rgba(color_str)
    if rgba and rgba[3] < 0.05:
        return None
    return {
        'pt': float(m.group(1)),
        'style': style,
        'color': css_color_to_hex(color_str) or '000000',
        'alpha': rgba[3] if rgba else 1.0,
    }


def _parse_border_data(s: dict) -> dict | None:
    """Returns uniform border data for shape.line, or None.
    Only returns data when all visible sides share the same color and thickness.
    Mixed borders are handled by _add_partial_borders instead."""
    side_results = {}
    for key in ('borderTop', 'borderRight', 'borderBottom', 'borderLeft'):
        v = s.get(key, '')
        if v and not v.startswith('0px'):
            parsed = _parse_border_one(v)
            if parsed:
                side_results[key] = parsed

    if side_results:
        # Check if all sides are uniform (same color and thickness)
        values = list(side_results.values())
        first = values[0]
        all_same = all(
            v['color'] == first['color'] and abs(v['pt'] - first['pt']) < 0.1
            for v in values
        )
        if all_same and len(side_results) >= 3:
            return first
        else:
            # Mixed colors/thickness or few sides → handled by _add_partial_borders
            return None

    border_str = s.get('border', '')
    if not border_str or border_str.startswith('0px') or border_str == 'none':
        return None
    return _parse_border_one(border_str)


def _parse_partial_borders(s: dict) -> dict:
    """Returns {side_name: border_data} for sides that have visible borders."""
    side_results = {}
    for key in ('borderTop', 'borderRight', 'borderBottom', 'borderLeft'):
        v = s.get(key, '')
        if v and not v.startswith('0px'):
            parsed = _parse_border_one(v)
            if parsed:
                side_results[key] = parsed
    return side_results


def _apply_gradient_fill_to_shape(shape, grad: dict):
    """Apply CSS gradient to a shape via direct spPr XML manipulation."""
    sp_pr = shape._element.spPr

    # Remove any existing fill elements
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:pattFill'):
        existing = sp_pr.find(qn(tag))
        if existing is not None:
            sp_pr.remove(existing)

    grad_fill = sp_pr.makeelement(qn('a:gradFill'), {})
    gs_lst = grad_fill.makeelement(qn('a:gsLst'), {})

    for stop in grad['stops']:
        gs = gs_lst.makeelement(qn('a:gs'), {'pos': str(int(stop['position'] * 1000))})
        rgb = css_color_to_rgb(stop['color'])
        if rgb:
            srgb_clr = gs.makeelement(qn('a:srgbClr'), {'val': f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'})
            gs.append(srgb_clr)
            gs_lst.append(gs)

    grad_fill.append(gs_lst)

    if grad['type'] == 'linear':
        ooxml_angle = int(((grad['angle'] + 270) % 360) * 60000)
        lin = grad_fill.makeelement(qn('a:lin'), {'ang': str(ooxml_angle), 'scaled': '0'})
        grad_fill.append(lin)

    # Insert after prstGeom
    prst_geom = sp_pr.find(qn('a:prstGeom'))
    if prst_geom is not None:
        prst_geom.addnext(grad_fill)
    else:
        sp_pr.append(grad_fill)


def _set_fill_transparency(shape, opacity_pct):
    """Set fill opacity via XML. opacity_pct: 100000=fully opaque, 0=fully transparent."""
    sp_pr = shape._element.spPr
    solid_fill = sp_pr.find(qn('a:solidFill'))
    if solid_fill is not None:
        srgb = solid_fill.find(qn('a:srgbClr'))
        if srgb is not None:
            for existing in srgb.findall(qn('a:alpha')):
                srgb.remove(existing)
            alpha_el = srgb.makeelement(qn('a:alpha'), {'val': str(opacity_pct)})
            srgb.append(alpha_el)


def _apply_border(shape, s: dict):
    border_data = _parse_border_data(s)
    if border_data:
        ln = shape.line
        ln.width = Pt(border_data['pt'])
        rgb = css_color_to_rgb(f'#{border_data["color"]}')
        if rgb:
            ln.color.rgb = RGBColor(*rgb)
        # Apply dash style
        _apply_line_dash(shape, border_data.get('style', 'solid'))
        alpha = border_data.get('alpha', 1.0)
        if alpha < 1.0:
            ln_elem = shape._element.spPr.find(qn('a:ln'))
            if ln_elem is not None:
                solid = ln_elem.find(qn('a:solidFill'))
                if solid is not None:
                    srgb = solid.find(qn('a:srgbClr'))
                    if srgb is not None:
                        alpha_el = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha * 100000))})
                        srgb.append(alpha_el)
    else:
        shape.line.fill.background()


def _apply_line_dash(shape, style: str):
    """Apply CSS border-style (dashed, dotted) to OOXML line dash."""
    if style in ('dashed', 'dotted'):
        ln_elem = shape._element.spPr.find(qn('a:ln'))
        if ln_elem is None:
            return
        dash_map = {
            'dashed': 'dash',
            'dotted': 'sysDot',
        }
        prst_dash = ln_elem.makeelement(qn('a:prstDash'), {'val': dash_map[style]})
        ln_elem.append(prst_dash)


def _add_partial_borders(slide, s: dict, x, y, w, h):
    """Add individual line shapes for sides that have non-uniform borders."""
    partial = _parse_partial_borders(s)
    uniform = _parse_border_data(s)
    if uniform or not partial:
        return  # uniform border handled by _apply_border, or no borders at all

    for side, data in partial.items():
        # data['pt'] is CSS px value of border width
        line_width = Pt(data['pt'])  # Pt handles conversion to EMU
        rgb = css_color_to_rgb(f'#{data["color"]}')
        if not rgb:
            continue

        # Calculate line endpoints based on side
        # Offset by half the line width so the line aligns with the box edge
        if side == 'borderTop':
            x1, y1, x2, y2 = x, y, x + w, y
        elif side == 'borderBottom':
            x1, y1, x2, y2 = x, y + h, x + w, y + h
        elif side == 'borderLeft':
            x1, y1, x2, y2 = x, y, x, y + h
        elif side == 'borderRight':
            x1, y1, x2, y2 = x + w, y, x + w, y + h
        else:
            continue

        connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT = 1
        _clear_theme_style(connector)
        connector.line.width = line_width
        connector.line.color.rgb = RGBColor(*rgb)

        # Apply dash style (dashed, dotted)
        style = data.get('style', 'solid')
        if style in ('dashed', 'dotted'):
            sp_pr = connector._element.find(qn('p:spPr'))
            if sp_pr is None:
                sp_pr = connector._element.spPr
            if sp_pr is not None:
                ln_elem = sp_pr.find(qn('a:ln'))
                if ln_elem is not None:
                    dash_map = {'dashed': 'dash', 'dotted': 'sysDot'}
                    prst_dash = ln_elem.makeelement(qn('a:prstDash'), {'val': dash_map[style]})
                    ln_elem.append(prst_dash)

        alpha = data.get('alpha', 1.0)
        if alpha < 1.0:
            ln_elem = connector._element.find(qn('a:ln'))
            if ln_elem is None:
                sp_pr = connector._element.spPr
                if sp_pr is not None:
                    ln_elem = sp_pr.find(qn('a:ln'))
            if ln_elem is not None:
                solid = ln_elem.find(qn('a:solidFill'))
                if solid is not None:
                    srgb = solid.find(qn('a:srgbClr'))
                    if srgb is not None:
                        alpha_el = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha * 100000))})
                        srgb.append(alpha_el)


def _parse_box_shadow(shadow: str):
    """Parse CSS box-shadow string and return (dx, dy, blur, rgba) or None.
    Filters out subtle/decorative shadows that shouldn't appear in PPT."""
    if not shadow or shadow == 'none':
        return None

    m = re.match(
        r'(?:rgba?\([^)]+\)\s+)?([-\d.]+)px\s+([-\d.]+)px\s+([-\d.]+)px(?:\s+([-\d.]+)px)?\s*(rgba?\([^)]+\))?',
        shadow.strip()
    )
    if not m:
        m = re.match(
            r'(rgba?\([^)]+\))\s+([-\d.]+)px\s+([-\d.]+)px\s+([-\d.]+)px(?:\s+([-\d.]+)px)?',
            shadow.strip()
        )
        if m:
            color_str = m.group(1)
            dx, dy, blur = float(m.group(2)), float(m.group(3)), float(m.group(4))
        else:
            return None
    else:
        dx, dy, blur = float(m.group(1)), float(m.group(2)), float(m.group(3))
        color_str = m.group(5) or ''
        if not color_str:
            cm = re.match(r'(rgba?\([^)]+\))', shadow.strip())
            color_str = cm.group(1) if cm else 'rgba(0,0,0,0.2)'

    rgba = css_color_to_rgba(color_str)
    if not rgba:
        rgba = (0, 0, 0, 0.2)

    # Filter out very subtle shadows:
    # - Very low alpha (< 0.15) → decorative/invisible
    # - Zero offset + zero blur → no visible effect
    # - Very small combined effect (alpha * blur < threshold)
    alpha = rgba[3]
    if alpha < 0.1:
        return None
    total_offset = abs(dx) + abs(dy)
    if total_offset < 0.5 and blur < 1:
        return None
    # Effective visibility: alpha * max(blur, offset) should be significant
    effective = alpha * max(blur, total_offset)
    if effective < 1.0:
        return None

    return (dx, dy, blur, rgba)


def _apply_shadow(shape, s: dict):
    """Apply CSS box-shadow to shape via OOXML effectLst."""
    parsed = _parse_box_shadow(s.get('boxShadow', ''))
    if not parsed:
        return
    dx, dy, blur, rgba = parsed

    sp = shape._element
    sp_pr = sp.spPr
    if sp_pr is None:
        return

    effect_lst = sp_pr.find(qn('a:effectLst'))
    if effect_lst is None:
        effect_lst = sp_pr.makeelement(qn('a:effectLst'), {})
        sp_pr.append(effect_lst)

    blur_emu = int(blur * 12700)
    dx_emu = int(dx * 12700)
    dy_emu = int(dy * 12700)
    dist = int(math.sqrt(dx_emu**2 + dy_emu**2))
    direction = int(math.atan2(dy_emu, dx_emu) * 60000 * 180 / math.pi) if dist > 0 else 0
    if direction < 0:
        direction += 360 * 60000

    outer_shdw = effect_lst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(blur_emu),
        'dist': str(dist),
        'dir': str(direction),
        'algn': 'tl',
        'rotWithShape': '0',
    })
    srgb = outer_shdw.makeelement(qn('a:srgbClr'), {
        'val': f'{rgba[0]:02X}{rgba[1]:02X}{rgba[2]:02X}'
    })
    if rgba[3] < 1.0:
        alpha_el = srgb.makeelement(qn('a:alpha'), {'val': str(int(rgba[3] * 100000))})
        srgb.append(alpha_el)
    outer_shdw.append(srgb)
    effect_lst.append(outer_shdw)


def _apply_text_shadow(shape, s: dict):
    """Apply CSS text-shadow to shape. Skip if box-shadow already added an effect."""
    shadow = s.get('textShadow', '')
    if not shadow or shadow == 'none':
        return

    # Don't add text-shadow if box-shadow already applied an effect (avoids double shadow)
    if _parse_box_shadow(s.get('boxShadow', '')):
        return

    # Parse: color dx dy blur  or  dx dy blur color
    m = re.match(
        r'(rgba?\([^)]+\))\s+([-\d.]+)px\s+([-\d.]+)px\s+([-\d.]+)px',
        shadow.strip()
    )
    if not m:
        m = re.match(
            r'([-\d.]+)px\s+([-\d.]+)px\s+([-\d.]+)px\s*(rgba?\([^)]+\))?',
            shadow.strip()
        )
        if not m:
            return
        dx, dy, blur = float(m.group(1)), float(m.group(2)), float(m.group(3))
        color_str = m.group(4) or 'rgba(0,0,0,0.3)'
    else:
        color_str = m.group(1)
        dx, dy, blur = float(m.group(2)), float(m.group(3)), float(m.group(4))

    rgba = css_color_to_rgba(color_str)
    if not rgba:
        return

    # Filter subtle text shadows (same logic as box-shadow)
    alpha = rgba[3]
    if alpha < 0.1:
        return
    total_offset = abs(dx) + abs(dy)
    if total_offset < 0.5 and blur < 1:
        return
    effective = alpha * max(blur, total_offset)
    if effective < 1.0:
        return

    sp = shape._element
    sp_pr = sp.spPr
    if sp_pr is None:
        return

    effect_lst = sp_pr.find(qn('a:effectLst'))
    if effect_lst is None:
        effect_lst = sp_pr.makeelement(qn('a:effectLst'), {})
        sp_pr.append(effect_lst)

    blur_emu = int(blur * 12700)
    dx_emu = int(dx * 12700)
    dy_emu = int(dy * 12700)
    dist = int(math.sqrt(dx_emu**2 + dy_emu**2))
    direction = int(math.atan2(dy_emu, dx_emu) * 60000 * 180 / math.pi) if dist > 0 else 0
    if direction < 0:
        direction += 360 * 60000

    outer_shdw = effect_lst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(blur_emu),
        'dist': str(dist),
        'dir': str(direction),
        'algn': 'tl',
        'rotWithShape': '0',
    })
    srgb = outer_shdw.makeelement(qn('a:srgbClr'), {
        'val': f'{rgba[0]:02X}{rgba[1]:02X}{rgba[2]:02X}'
    })
    if rgba[3] < 1.0:
        alpha_el = srgb.makeelement(qn('a:alpha'), {'val': str(int(rgba[3] * 100000))})
        srgb.append(alpha_el)
    outer_shdw.append(srgb)
    effect_lst.append(outer_shdw)


def _apply_opacity(shape, s: dict):
    opacity = s.get('opacity', '1')
    if opacity and opacity != '1':
        try:
            val = float(opacity)
            if val < 1:
                # Set transparency via XML (not directly supported in python-pptx API)
                sp = shape._element
                sp_pr = sp.find(qn('p:spPr'))
                if sp_pr is None:
                    sp_pr = sp.find(qn('xdr:spPr'))
                # Opacity not trivially supported - skip for now
        except (ValueError, TypeError):
            pass


def _apply_round_corners(shape, border_radius: str):
    try:
        # Parse multi-value borderRadius: "8px 8px 0px 0px" or single "8px"
        parts = border_radius.strip().split()
        px_values = []
        for part in parts:
            clean = part.replace('px', '').replace('%', '').replace(',', '')
            try:
                px_values.append(float(clean))
            except ValueError:
                continue

        if not px_values:
            return

        # Use max corner radius for roundRect (OOXML only supports uniform rounding)
        px = max(px_values)
        if px < 1:
            return

        is_percent = '%' in border_radius
        if is_percent:
            min_dim = min(shape.width, shape.height)
            emu = int(min_dim * px / 100)
        else:
            emu = int(px * PX_TO_EMU)

        sp_pr = shape._element.spPr
        prst_geom = sp_pr.find(qn('a:prstGeom'))
        if prst_geom is not None:
            prst_geom.set('prst', 'roundRect')
            av_lst = prst_geom.find(qn('a:avLst'))
            if av_lst is None:
                av_lst = prst_geom.makeelement(qn('a:avLst'), {})
                prst_geom.append(av_lst)
            else:
                for child in list(av_lst):
                    av_lst.remove(child)
            min_dim = min(shape.width, shape.height)
            if min_dim > 0:
                adj_val = min(int((emu / min_dim) * 50000), 50000)
                gd = av_lst.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': f'val {adj_val}'})
                av_lst.append(gd)
    except (ValueError, TypeError):
        pass


def _set_anchor(txbox, anchor: str):
    body_props = txbox.text_frame._txBody.find(qn('a:bodyPr'))
    if body_props is not None:
        mapping = {'top': 't', 'middle': 'ctr', 'bottom': 'b'}
        body_props.set('anchor', mapping.get(anchor, 't'))


def _set_margins(txbox, padding: str):
    body_props = txbox.text_frame._txBody.find(qn('a:bodyPr'))
    if body_props is None:
        return

    parts = padding.split()
    px_values = []
    for p in parts:
        try:
            px_values.append(float(p.replace('px', '')))
        except ValueError:
            px_values.append(0)

    if len(px_values) == 1:
        t = r = b = l = px_values[0]
    elif len(px_values) == 2:
        t = b = px_values[0]
        r = l = px_values[1]
    elif len(px_values) == 3:
        t, r, b = px_values[0], px_values[1], px_values[2]
        l = r
    elif len(px_values) >= 4:
        t, r, b, l = px_values[0], px_values[1], px_values[2], px_values[3]
    else:
        t = r = b = l = 0

    body_props.set('lIns', str(int(l * PX_TO_EMU)))
    body_props.set('rIns', str(int(r * PX_TO_EMU)))
    body_props.set('tIns', str(int(t * PX_TO_EMU)))
    body_props.set('bIns', str(int(b * PX_TO_EMU)))
